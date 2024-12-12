import os
from psd_tools import PSDImage
from pptx import Presentation
from pptx.util import Inches

# Fungsi untuk konversi PSD ke PNG
def convert_psd_to_png(psd_path, output_folder):
    psd = PSDImage.open(psd_path)
    png_image_path = os.path.join(output_folder, os.path.basename(psd_path).replace('.psd', '.png'))
    psd.composite().save(png_image_path, 'PNG')
    return png_image_path

# Fungsi untuk membuat PPT dengan gambar-gambar
def create_ppt_with_images(image_paths, ppt_output_path):
    prs = Presentation()
    for image_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(6))
    prs.save(ppt_output_path)

# Fungsi utama automasi
def psd_to_ppt_automation(psd_folder, ppt_output_path):
    image_paths = []
    for filename in os.listdir(psd_folder):
        if filename.endswith('.psd'):
            # Konversi setiap file PSD ke PNG
            png_path = convert_psd_to_png(os.path.join(psd_folder, filename), psd_folder)
            image_paths.append(png_path)
    # Buat PPT dengan semua gambar PNG yang dihasilkan
    create_ppt_with_images(image_paths, ppt_output_path)

# Memastikan kode ini berjalan saat dipanggil
if __name__ == "__main__":
    psd_folder = '/Documents/TUGAS KULIAH/BAHASA INDONESIA SMS 1/ppt'
